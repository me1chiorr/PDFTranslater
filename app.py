import streamlit as st
import tiktoken
from langchain.text_splitter import RecursiveCharacterTextSplitter
import secrets
import pptx
import openai
from io import BufferedReader
from fastapi import UploadFile
import mimetypes
from PyPDF2 import PdfReader
import hashlib
import docx2txt
import csv
from pydantic import BaseModel
from typing import List, Optional
import os
from typing import List
from dotenv import load_dotenv, find_dotenv

class Document(BaseModel):
    id: Optional[str] = None
    text: str


tokenizer = tiktoken.get_encoding('cl100k_base')


# create the length function
def tiktoken_len(text):
    tokens = tokenizer.encode(
        text,
        disallowed_special=()
    )
    return len(tokens)


def text_splitter(chunk_size=2048):
    text_splitter = RecursiveCharacterTextSplitter(
        # Set a really small chunk size, just to show.
        chunk_size=chunk_size,
        chunk_overlap=5,
        length_function=tiktoken_len,
    )
    return text_splitter


# 2
_ = load_dotenv(find_dotenv())  # read local .env file
file_folder = os.environ['FILE_FOLDER'] if 'FILE_FOLDER' in os.environ else "translated_file"


async def save_as_txt(content_list):
    token = secrets.token_hex(16)
    if not os.path.exists(file_folder):
        os.makedirs(file_folder)

    with open(f"{file_folder}/{token}.txt", "w+") as f:
        for i in range(len(content_list)):
            f.write(content_list[i].translated_content)

    return token


# 3


async def get_document_from_file(file: UploadFile) -> Document:
    extracted_text = await extract_text_from_form_file(file)
    doc = Document(text=extracted_text)

    return doc


def extract_text_from_filepath(filepath: str, mimetype: Optional[str] = None) -> str:
    """Return the text content of a file given its filepath."""

    if mimetype is None:
        # Get the mimetype of the file based on its extension
        mimetype, _ = mimetypes.guess_type(filepath)

    if not mimetype:
        if filepath.endswith(".md"):
            mimetype = "text/markdown"
        else:
            raise Exception("Unsupported file type")

    # Open the file in binary mode
    file = open(filepath, "rb")
    extracted_text = extract_text_from_file(file, mimetype)

    return extracted_text


def extract_text_from_file(file: BufferedReader, mimetype: str) -> str:
    if mimetype == "application/pdf":
        # Extract text from pdf using PyPDF2
        reader = PdfReader(file)
        extracted_text = ""
        for page in reader.pages:
            extracted_text += page.extract_text()
    elif mimetype == "text/plain" or mimetype == "text/markdown":
        # Read text from plain text file
        extracted_text = file.read().decode("utf-8")
    elif (
            mimetype
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        # Extract text from docx using docx2txt
        extracted_text = docx2txt.process(file)
    elif mimetype == "text/csv":
        # Extract text from csv using csv module
        extracted_text = ""
        decoded_buffer = (line.decode("utf-8") for line in file)
        reader = csv.reader(decoded_buffer)
        for row in reader:
            extracted_text += " ".join(row) + "\n"
    elif (
            mimetype
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
    else:
        # Unsupported file type
        file.close()
        raise ValueError("Unsupported file type: {}".format(mimetype))

    file.close()
    return extracted_text


# Extract text from a file based on its mimetype
async def extract_text_from_form_file(file: UploadFile):
    """Return the text content of a file."""
    # get the file body from the upload file object
    mimetype = file.content_type
    print(f"mimetype: {mimetype}")

    file_stream = await file.read()

    hash_code = hashlib.sha256(file_stream).hexdigest()

    if not os.path.exists("./temp_files/"):
        os.makedirs("./temp_files/")
        print("Temporary Folder created successfully!")

    temp_file_path = f"./temp_files/{hash_code}"

    try:
        with open(temp_file_path, "wb") as f:
            f.write(file_stream)
        extracted_text = extract_text_from_filepath(temp_file_path, mimetype)

    except Exception as e:
        raise e

    os.remove(temp_file_path)

    return extracted_text


# 4
_ = load_dotenv(find_dotenv())  # read local .env file


def call_openai(sys_prompt: str, user_prompt, api_type: str):
    messages = [
        {"role": "system", "content": sys_prompt}
    ]

    if api_type == ApiType.azure.value:
        engine = os.environ['AZURE_DEPLOYMENT_NAME']
        openai.api_type = api_type
        openai.api_key = os.environ['AZURE_API_KEY']
        openai.api_base = os.environ['AZURE_API_BASE']
        openai.api_version = os.environ['AZURE_API_VERSION']
        response = openai.ChatCompletion.create(
            engine=engine,
            messages=messages,
            temperature=float(os.environ['TEMPERATURE'])
        )
        return response


    elif api_type == ApiType.open_ai.value:
        openai.api_key = st.secrets["OPENAI_API_KEY"]
        for text in user_prompt:
            messages.append({
                "role": "user", "content": text.text
            })
        response = openai.ChatCompletion.create(
            model=os.environ["OPENAI_MODEL"],
            messages=messages,
            temperature=float(os.environ['TEMPERATURE'])
        )
        return response


class TranslateResult(BaseModel):
    original_content: str
    translated_content: str


class TranslateResponse(BaseModel):
    results: List[TranslateResult]


class TranslatedFileResponse(BaseModel):
    result: str


from pydantic import BaseModel
from typing import List, Optional
from enum import Enum


class ApiType(str, Enum):
    open_ai = 'open_ai'
    azure = 'azure'


class TranslateType(str, Enum):
    ZH_EN = 'zh_en'
    EN_ZH = 'en_zh'


_ = load_dotenv(find_dotenv())  # read local .env file


def get_translate_results(texts: List[Document], translate_type: str, api_type: str) -> List[TranslateResult]:
    results = []

    sys_prompt = os.environ[translate_type.upper()]
    response = call_openai(sys_prompt=sys_prompt, user_prompt=texts, api_type=api_type)
    if response:
        choices = response["choices"]
        completion = choices[0].message.content.strip()
        return completion

    return results


def main():
    st.title("PDF Translator")

    file = st.file_uploader("Upload PDF file", type="pdf")



    add_selectbox = st.sidebar.selectbox(
        "How would you like to be contacted?",
        ("Email", "Home phone", "Mobile phone")
    )

    add_toggle = st.sidebar.toggle(
        "Want to Stay Turned?",

    )

    if file is not None:
        pdf_reader = PdfReader(file)

        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        splitter = text_splitter()

        chunks = splitter.split_text(text)

        documents = [Document(text=chunk) for chunk in chunks]
        print(documents)

        result = get_translate_results(documents, translate_type='ru_kz', api_type='open_ai')




        st.text_area(
            'Переведенный текст',
            result
        )


main()
